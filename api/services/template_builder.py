"""
Template Builder — Creates professional themed PPTX presentations.

Features:
  - Multiple theme presets (dark purple, ocean, forest, sunset, midnight, crimson)
  - Gradient backgrounds per theme
  - Styled typography (Calibri Light titles, Calibri body)
  - Decorative accent bar at slide bottom
  - Auto-fit font sizing so content never overflows
  - Image placement support
"""
import logging
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE

logger = logging.getLogger("odin_api.services.template_builder")

# ── Theme Presets ───────────────────────────────────────────
THEMES = {
    "dark_purple": {
        "bg_dark":       RGBColor(0x0F, 0x0A, 0x1A),
        "bg_gradient":   RGBColor(0x1A, 0x0A, 0x2E),
        "title":         RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle":      RGBColor(0xC4, 0xB5, 0xFD),
        "body":          RGBColor(0xE0, 0xE0, 0xE0),
        "accent":        RGBColor(0x7C, 0x3A, 0xED),
        "accent_light":  RGBColor(0xA7, 0x8B, 0xFA),
        "muted":         RGBColor(0x9C, 0xA3, 0xAF),
    },
    "ocean": {
        "bg_dark":       RGBColor(0x03, 0x13, 0x1A),
        "bg_gradient":   RGBColor(0x06, 0x1E, 0x33),
        "title":         RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle":      RGBColor(0x7D, 0xD3, 0xFC),
        "body":          RGBColor(0xD6, 0xED, 0xF5),
        "accent":        RGBColor(0x06, 0x8F, 0xCF),
        "accent_light":  RGBColor(0x38, 0xBD, 0xF8),
        "muted":         RGBColor(0x8E, 0xA8, 0xBB),
    },
    "forest": {
        "bg_dark":       RGBColor(0x07, 0x15, 0x0B),
        "bg_gradient":   RGBColor(0x0A, 0x28, 0x14),
        "title":         RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle":      RGBColor(0x86, 0xEF, 0xAC),
        "body":          RGBColor(0xD8, 0xF0, 0xDB),
        "accent":        RGBColor(0x16, 0xA3, 0x4A),
        "accent_light":  RGBColor(0x4A, 0xDE, 0x80),
        "muted":         RGBColor(0x8C, 0xAF, 0x94),
    },
    "sunset": {
        "bg_dark":       RGBColor(0x1A, 0x0A, 0x05),
        "bg_gradient":   RGBColor(0x30, 0x10, 0x08),
        "title":         RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle":      RGBColor(0xFD, 0xBA, 0x74),
        "body":          RGBColor(0xF0, 0xE0, 0xD0),
        "accent":        RGBColor(0xEA, 0x58, 0x0C),
        "accent_light":  RGBColor(0xFB, 0x92, 0x3C),
        "muted":         RGBColor(0xBB, 0xA0, 0x8C),
    },
    "midnight": {
        "bg_dark":       RGBColor(0x0A, 0x0A, 0x14),
        "bg_gradient":   RGBColor(0x12, 0x12, 0x22),
        "title":         RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle":      RGBColor(0x93, 0xC5, 0xFD),
        "body":          RGBColor(0xD4, 0xDE, 0xEC),
        "accent":        RGBColor(0x25, 0x63, 0xEB),
        "accent_light":  RGBColor(0x60, 0xA5, 0xFA),
        "muted":         RGBColor(0x88, 0x99, 0xAA),
    },
    "crimson": {
        "bg_dark":       RGBColor(0x1A, 0x06, 0x08),
        "bg_gradient":   RGBColor(0x2D, 0x0A, 0x0F),
        "title":         RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle":      RGBColor(0xFC, 0xA5, 0xA5),
        "body":          RGBColor(0xF0, 0xDC, 0xDC),
        "accent":        RGBColor(0xDC, 0x26, 0x26),
        "accent_light":  RGBColor(0xF8, 0x71, 0x71),
        "muted":         RGBColor(0xBB, 0x8C, 0x8C),
    },
    "emerald_gold": {
        "bg_dark":       RGBColor(0x0B, 0x14, 0x10),
        "bg_gradient":   RGBColor(0x12, 0x24, 0x1A),
        "title":         RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle":      RGBColor(0xFD, 0xE6, 0x8A),
        "body":          RGBColor(0xE8, 0xED, 0xE0),
        "accent":        RGBColor(0xCA, 0x88, 0x10),
        "accent_light":  RGBColor(0xFA, 0xCC, 0x15),
        "muted":         RGBColor(0xA0, 0xAA, 0x90),
    },
    "rose": {
        "bg_dark":       RGBColor(0x18, 0x08, 0x14),
        "bg_gradient":   RGBColor(0x28, 0x0C, 0x22),
        "title":         RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle":      RGBColor(0xF9, 0xA8, 0xD4),
        "body":          RGBColor(0xF0, 0xDC, 0xE8),
        "accent":        RGBColor(0xDB, 0x27, 0x77),
        "accent_light":  RGBColor(0xF4, 0x72, 0xB6),
        "muted":         RGBColor(0xB0, 0x8C, 0xA0),
    },
}

AVAILABLE_THEMES = list(THEMES.keys())
DEFAULT_THEME = "dark_purple"

# ── Slide Dimensions (16:9) ─────────────────────────────────
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ── Font Configuration ──────────────────────────────────────
FONT_TITLE = "Calibri Light"
FONT_BODY = "Calibri"

TITLE_SLIDE_TITLE_SIZE = Pt(40)
TITLE_SLIDE_SUBTITLE_SIZE = Pt(20)

CONTENT_TITLE_SIZE = Pt(28)
CONTENT_BODY_SIZE = Pt(16)
CONTENT_BODY_MIN_SIZE = Pt(10)

ACCENT_BAR_HEIGHT = Inches(0.08)


def get_theme(theme_name: str | None = None) -> dict:
    """Get a theme palette by name. Falls back to default if not found."""
    if theme_name and theme_name in THEMES:
        return THEMES[theme_name]
    # Try fuzzy match
    if theme_name:
        name = theme_name.lower().strip().replace(" ", "_").replace("-", "_")
        if name in THEMES:
            return THEMES[name]
        # Partial match
        for key in THEMES:
            if name in key or key in name:
                return THEMES[key]
    return THEMES[DEFAULT_THEME]


def _set_slide_gradient(slide, colors):
    """Apply dark gradient background to a slide."""
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_angle = 315

    stop0 = fill.gradient_stops[0]
    stop0.position = 0.0
    stop0.color.rgb = colors["bg_dark"]

    stop1 = fill.gradient_stops[1]
    stop1.position = 1.0
    stop1.color.rgb = colors["bg_gradient"]


def _add_accent_bar(slide, colors, y_position=None):
    """Add decorative accent gradient bar at the bottom."""
    if y_position is None:
        y_position = SLIDE_HEIGHT - ACCENT_BAR_HEIGHT - Inches(0.3)

    bar = slide.shapes.add_shape(
        1, Inches(0.5), y_position,
        SLIDE_WIDTH - Inches(1.0), ACCENT_BAR_HEIGHT,
    )
    bar.line.fill.background()

    fill = bar.fill
    fill.gradient()
    fill.gradient_angle = 0
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[0].color.rgb = colors["accent"]
    fill.gradient_stops[1].position = 1.0
    fill.gradient_stops[1].color.rgb = colors["accent_light"]


def _add_decorative_shape(slide, x, y, w, h, color, shape_type=1):
    """Add a small solid-color decorative shape (rectangle or oval)."""
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _estimate_text_lines(text, chars_per_line=70):
    """Estimate line count for text block."""
    total = 0
    for line in text.split("\n"):
        if not line.strip():
            total += 1
        else:
            total += max(1, len(line) // chars_per_line + 1)
    return total


def _calculate_font_size(content_text, available_height_inches,
                         max_size=CONTENT_BODY_SIZE,
                         min_size=CONTENT_BODY_MIN_SIZE):
    """Auto-shrink font size so content fits within available height."""
    LINE_HEIGHT_RATIO = 0.028  # inches per pt per line

    current_pt = max_size.pt if hasattr(max_size, "pt") else max_size
    min_pt = min_size.pt if hasattr(min_size, "pt") else min_size

    while current_pt >= min_pt:
        chars_per_line = int(80 * (16 / current_pt))
        estimated_lines = _estimate_text_lines(content_text, chars_per_line)
        total_height = estimated_lines * current_pt * LINE_HEIGHT_RATIO

        if total_height <= available_height_inches:
            return Pt(current_pt)
        current_pt -= 1

    return Pt(min_pt)


def _strip_html_and_markdown(text):
    """Remove HTML tags and markdown formatting from text."""
    text = re.sub(r'<[^>]+>', '', text)
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)
    text = re.sub(r'^#{1,6}\s+', '', text, flags=re.MULTILINE)
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


def _format_content_text(text_frame, content, font_size, colors):
    """Format content with bullet points and styled paragraphs."""
    text_frame.clear()
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.word_wrap = True

    content = _strip_html_and_markdown(content)

    lines = content.split("\n")
    first = True

    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue

        para = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
        first = False

        is_bullet = stripped.startswith(("- ", "• ", "* ", "– "))
        if is_bullet:
            stripped = stripped.lstrip("-•*– ").strip()
            run = para.add_run()
            run.text = f"  •  {stripped}"
            para.space_before = Pt(2)
            para.space_after = Pt(2)
        else:
            run = para.add_run()
            run.text = stripped
            para.space_before = Pt(4)
            para.space_after = Pt(4)

        run.font.name = FONT_BODY
        run.font.size = font_size
        run.font.color.rgb = colors["body"]


def build_title_slide(prs, title, subtitle="", colors=None):
    """Create a visually striking title slide."""
    if colors is None:
        colors = THEMES[DEFAULT_THEME]

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    _set_slide_gradient(slide, colors)

    # Small accent decoration top-left
    _add_decorative_shape(slide, Inches(0.3), Inches(0.3),
                          Inches(0.5), Pt(4), colors["accent"])

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(2.0), Inches(11.3), Inches(2.0)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.name = FONT_TITLE
    run.font.size = TITLE_SLIDE_TITLE_SIZE
    run.font.color.rgb = colors["title"]
    run.font.bold = True

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(2.0), Inches(4.2), Inches(9.3), Inches(1.0)
        )
        tf2 = sub_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.name = FONT_BODY
        run2.font.size = TITLE_SLIDE_SUBTITLE_SIZE
        run2.font.color.rgb = colors["subtitle"]

    _add_accent_bar(slide, colors, Inches(5.5))
    return slide


def build_content_slide(prs, title, content, image_path=None,
                        slide_number=None, colors=None):
    """
    Create a content slide with title, body, optional image.
    Auto-adjusts font size to fit content.
    """
    if colors is None:
        colors = THEMES[DEFAULT_THEME]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_gradient(slide, colors)

    # Slide number badge (top-right)
    if slide_number is not None:
        num_box = slide.shapes.add_textbox(
            SLIDE_WIDTH - Inches(1.2), Inches(0.3), Inches(0.8), Inches(0.4)
        )
        p_num = num_box.text_frame.paragraphs[0]
        p_num.alignment = PP_ALIGN.RIGHT
        run_num = p_num.add_run()
        run_num.text = str(slide_number)
        run_num.font.name = FONT_BODY
        run_num.font.size = Pt(12)
        run_num.font.color.rgb = colors["muted"]

    # If image, split layout 60/40
    has_image = image_path and Path(image_path).exists()
    if has_image:
        content_width = Inches(7.5)
        content_left = Inches(0.8)
        img_left = Inches(8.8)
        img_width = Inches(4.0)
        img_top = Inches(1.8)
        img_height = Inches(4.5)
    else:
        content_width = Inches(11.5)
        content_left = Inches(0.8)

    # Title
    title_box = slide.shapes.add_textbox(
        content_left, Inches(0.5), content_width, Inches(1.0)
    )
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.alignment = PP_ALIGN.LEFT
    run_title = p_title.add_run()
    run_title.text = title
    run_title.font.name = FONT_TITLE
    run_title.font.size = CONTENT_TITLE_SIZE
    run_title.font.color.rgb = colors["title"]
    run_title.font.bold = True

    # Title underline accent
    underline = slide.shapes.add_shape(
        1, content_left, Inches(1.45), Inches(2.0), Pt(3),
    )
    underline.line.fill.background()
    underline.fill.solid()
    underline.fill.fore_color.rgb = colors["accent"]

    # Body content
    body_top = Inches(1.7)
    body_height = 5.0
    body_box = slide.shapes.add_textbox(
        content_left, body_top, content_width, Inches(body_height)
    )
    tf_body = body_box.text_frame
    tf_body.word_wrap = True

    font_size = _calculate_font_size(content, body_height)
    logger.debug(f"Slide '{title}': auto-fit font size = {font_size.pt}pt")
    _format_content_text(tf_body, content, font_size, colors)

    # Add image if available
    if has_image:
        try:
            slide.shapes.add_picture(
                str(image_path), img_left, img_top, img_width, img_height
            )
            logger.info(f"Added image to slide: {image_path}")
        except Exception as e:
            logger.warning(f"Failed to add image: {e}")

    _add_accent_bar(slide, colors)
    return slide


def build_themed_presentation(slides_data=None, image_paths=None,
                              theme_name=None):
    """
    Build a complete themed presentation.

    Args:
        slides_data: List of dicts with: slide_number, title, content, narration
        image_paths: Dict mapping slide_number -> image file path
        theme_name: Theme preset name (e.g. 'ocean', 'forest', 'sunset')

    Returns:
        Presentation object
    """
    colors = get_theme(theme_name)
    used_theme = theme_name or DEFAULT_THEME
    logger.info(f"Building presentation with theme: {used_theme}")

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    if not slides_data:
        return prs

    if image_paths is None:
        image_paths = {}

    # First slide as title slide
    first = slides_data[0]
    subtitle = first.get("content", "")
    if len(subtitle) > 120:
        subtitle = subtitle[:120] + "..."
    build_title_slide(prs, title=first.get("title", "Presentation"),
                      subtitle=subtitle, colors=colors)

    # Remaining slides as content
    for i, sd in enumerate(slides_data[1:], start=2):
        img = image_paths.get(sd.get("slide_number"), None)
        build_content_slide(
            prs,
            title=sd.get("title", f"Slide {i}"),
            content=sd.get("content", ""),
            image_path=img,
            slide_number=i,
            colors=colors,
        )

    return prs
