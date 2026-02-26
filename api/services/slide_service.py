"""
Slide Service — PPTX creation, reading, and merge logic.
Extracted and refactored from odin_slides/presentation.py for API usage.
"""
import difflib
import logging
import os
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..core.config import settings

logger = logging.getLogger("odin_api.slides")


# ── Layout Helpers ──────────────────────────────────────────

def _find_most_similar_layout(prs: Presentation, target_name: str):
    """Find the most similar layout in the presentation to the target name."""
    layout_names = [layout.name for layout in prs.slide_layouts]
    closest_matches = difflib.get_close_matches(target_name, layout_names)
    if closest_matches:
        for layout in prs.slide_layouts:
            if layout.name == closest_matches[0]:
                return layout
    return None


def _find_content_placeholder(slide):
    """Find the content placeholder (idx=1) in a slide."""
    for shape in slide.shapes:
        if (shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER
                and shape.placeholder_format.idx == 1):
            return shape
    return None


# ── PPTX Operations ────────────────────────────────────────

def get_template_path(template_name: str | None = None) -> Path:
    """Resolve the path to a template file."""
    if template_name:
        path = settings.TEMPLATES_DIR / template_name
        if path.exists():
            return path
    # Look for any pptx in templates dir
    for f in settings.TEMPLATES_DIR.glob("*.pptx"):
        return f
    # Create a minimal default template
    return _create_default_template()


def _create_default_template() -> Path:
    """Create a minimal default PPTX template."""
    prs = Presentation()
    path = settings.TEMPLATES_DIR / "default_template.pptx"
    prs.save(str(path))
    logger.info(f"Created default template at {path}")
    return path


async def create_pptx(
    slides: list[dict],
    template_path: Path | str | None = None,
    output_path: Path | str | None = None,
    theme_name: str | None = None,
) -> Path:
    """
    Create a themed PPTX file from slide data with images.

    Uses template_builder for professional gradient theme and
    image_service for relevant stock photos.

    Args:
        slides: List of slide dicts with keys: slide_number, title, content, narration, image_keyword.
        template_path: Ignored (kept for backward compatibility). Uses themed builder.
        output_path: Where to save. Uses temp file if None.
        theme_name: Theme preset name (e.g. 'ocean', 'forest', 'sunset').

    Returns:
        Path to the created PPTX file.
    """
    from .template_builder import build_themed_presentation
    from .image_service import fetch_images_for_slides

    # Fetch images for all slides (graceful: returns {} if no API key)
    image_paths = await fetch_images_for_slides(slides)

    # Build themed presentation
    prs = build_themed_presentation(slides_data=slides, image_paths=image_paths,
                                    theme_name=theme_name)

    # Determine output path
    if output_path is None:
        output_path = settings.TEMP_DIR / f"presentation_{id(prs)}.pptx"
    output_path = Path(output_path)

    prs.save(str(output_path))
    logger.info(f"Themed presentation saved to {output_path} ({len(slides)} slides)")
    return output_path


def read_pptx(pptx_path: Path | str) -> list[dict]:
    """
    Read slide data from a PPTX file.

    Args:
        pptx_path: Path to the PPTX file.

    Returns:
        List of slide dicts.
    """
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        return []

    prs = Presentation(str(pptx_path))
    slides = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_info = {
            "slide_number": float(slide_number),
            "title": "",
            "content": "",
            "narration": "",
        }

        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape == slide.shapes[0]:
                    slide_info["title"] = shape.text
                else:
                    slide_info["content"] += shape.text + "\n"

        # Extract notes
        notes_slide = slide.notes_slide
        for shape in notes_slide.shapes:
            if shape.has_text_frame:
                slide_info["narration"] += shape.text + "\n"

        slides.append(slide_info)

    return slides


# ── Merge Logic ─────────────────────────────────────────────

def merge_slides(existing_slides: list[dict], new_slides: list[dict]) -> list[dict]:
    """
    Merge new slides into existing slide deck.

    Rules (preserved from original odin-slides):
    - Same slide_number: replace existing with new
    - Negative slide_number: delete the corresponding slide
    - Decimal slide_number (e.g., 2.1): insert new slide
    - Final result: sorted by slide_number, then renumbered from 1

    Args:
        existing_slides: Current slide deck.
        new_slides: Slides from LLM response.

    Returns:
        Merged and renumbered slide list.
    """
    # Format slide numbers as strings for comparison
    for slide in existing_slides:
        slide["slide_number"] = "{:.1f}".format(float(slide["slide_number"]))
    for slide in new_slides:
        slide["slide_number"] = "{:.1f}".format(float(slide["slide_number"]))

    # Remove existing slides that have the same number as new slides
    new_numbers = set(s["slide_number"] for s in new_slides)
    existing_slides = [s for s in existing_slides if s["slide_number"] not in new_numbers]

    # Compute opposite numbers for deletion logic
    existing_numbers = set(s["slide_number"] for s in existing_slides)
    opposite_existing = set("{:.1f}".format(-float(n)) for n in existing_numbers)

    new_numbers_set = set(s["slide_number"] for s in new_slides)
    opposite_new = set("{:.1f}".format(-float(n)) for n in new_numbers_set)

    # Filter out slides marked for deletion
    filtered_existing = [s for s in existing_slides if s["slide_number"] not in opposite_new]
    filtered_new = [s for s in new_slides if s["slide_number"] not in opposite_existing]

    # Merge and sort
    merged = sorted(
        filtered_existing + filtered_new,
        key=lambda x: float(x["slide_number"])
    )

    # Renumber starting from 1
    for i, slide in enumerate(merged, start=1):
        slide["slide_number"] = i

    return merged


def slides_to_preview(slides: list[dict]) -> list[dict]:
    """
    Convert slide data to a format suitable for frontend preview.

    Args:
        slides: List of slide dicts.

    Returns:
        Cleaned slide data list for JSON response.
    """
    preview = []
    for slide in slides:
        preview.append({
            "slide_number": int(slide.get("slide_number", 0)),
            "title": slide.get("title", ""),
            "content": slide.get("content", ""),
            "narration": slide.get("narration", ""),
            "image_keyword": slide.get("image_keyword", ""),
        })
    return preview
